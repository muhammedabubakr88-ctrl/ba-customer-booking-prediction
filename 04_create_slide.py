"""
Task 2 - Step 4: Build the single-slide PowerPoint summary
"""
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Load our saved results
with open('outputs/results.json', 'r') as f:
    results = json.load(f)

test = results['test_metrics']
cv = results['cv_summary']

# BA-ish brand colors
NAVY = RGBColor(0x07, 0x2A, 0x5E)
RED = RGBColor(0xEF, 0x33, 0x40)
GREY = RGBColor(0x4D, 0x4D, 0x4D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

# --- Title bar ---
title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Predicting Customer Booking Completion"
run.font.size = Pt(30)
run.font.bold = True
run.font.color.rgb = NAVY

subtitle_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.9), Inches(12.5), Inches(0.4))
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Random Forest model \u2014 identifying customers likely to complete a booking, before departure"
run.font.size = Pt(14)
run.font.italic = True
run.font.color.rgb = GREY

# --- Chart image (left side) ---
slide.shapes.add_picture('outputs/feature_importance_chart.png',
                          Inches(0.4), Inches(1.5), height=Inches(5.6))

# --- Right-side metrics panel ---
panel_left = Inches(7.6)
panel_top = Inches(1.5)
panel_width = Inches(5.3)

box = slide.shapes.add_textbox(panel_left, panel_top, panel_width, Inches(0.4))
tf = box.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Model Performance"
run.font.size = Pt(18)
run.font.bold = True
run.font.color.rgb = NAVY

metrics_text = [
    ("ROC-AUC", f"{test['roc_auc']:.2f}", "Model distinguishes bookers from non-bookers well above chance (0.50)"),
    ("Recall", f"{test['recall']:.0%}", "Of customers who go on to book, the model correctly flags most of them"),
    ("Precision", f"{test['precision']:.0%}", "Of those flagged, about 3 in 10 actually book \u2014 expected given bookings are rare (15% of all customers)"),
]

top = Inches(2.0)
for label, val, desc in metrics_text:
    box = slide.shapes.add_textbox(panel_left, top, panel_width, Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{label}: {val}"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RED

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = desc
    run2.font.size = Pt(11)
    run2.font.color.rgb = GREY

    top += Inches(1.05)

# --- Key takeaway box at bottom right ---
takeaway_box = slide.shapes.add_textbox(panel_left, Inches(5.3), panel_width, Inches(2.0))
tf = takeaway_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Recommendation"
run.font.size = Pt(16)
run.font.bold = True
run.font.color.rgb = NAVY

p2 = tf.add_paragraph()
run2 = p2.add_run()
run2.text = (
    "Booking origin, route popularity, and trip length are the strongest predictors of "
    "booking completion. These signals can power a proactive targeting model \u2014 flagging "
    "high-likelihood customers for outreach before they reach the airport, rather than "
    "reacting after the fact."
)
run2.font.size = Pt(12)
run2.font.color.rgb = GREY

prs.save('outputs/BA_booking_prediction_summary.pptx')
print("Saved outputs/BA_booking_prediction_summary.pptx")